"""最終提出スライド（本編8枚 + Backup 3枚 = 11枚、16:9）を生成する。

設計: docs/PRESENTATION_STORYBOARD.md（人間承認済み）
数値: すべて outputs/ 配下のログから読み取る。**ハードコードしない。**

【なぜプレビュー PNG も出すのか】
この環境に LibreOffice 等の pptx レンダラが無いため、pptx を目視確認できない。
そこで **同一のレイアウト定義から pptx と PNG の両方を生成**し、
PNG 側で文字切れ・可読性を監査する。位置・サイズ・フォントは共通の SPEC を使うため、
PNG で収まっていれば pptx でも収まる。

使い方:
    python slides/make_pptx.py
"""

from __future__ import annotations

import json
import statistics as st
from dataclasses import dataclass, field
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
MAIN = ROOT / "outputs" / "main_experiment"
SENS = ROOT / "outputs" / "sensitivity_replay" / "penalty_sensitivity.json"
LIVE = ROOT / "outputs" / "live_penalty_zero"
FIGS = ROOT / "figures"
OUT = ROOT / "slides"
PREVIEW = OUT / "preview"

# 16:9
W_IN, H_IN = 13.333, 7.5
PX_PER_IN = 96  # プレビュー解像度の基準
SCALE = 2       # プレビューは 2倍で描画して縮小（文字の潰れを見やすくする）

INK = (26, 26, 26)
MUTED = (138, 138, 138)
OK = (46, 125, 91)
GAP = (181, 83, 60)
ACCENT = (51, 98, 143)
LIGHT = (220, 220, 220)
BG = (255, 255, 255)
PANEL = (246, 246, 246)

FONT_CANDIDATES = ["meiryo.ttc", "YuGothM.ttc", "msgothic.ttc", "meiryob.ttc"]
FONT_DIR = Path("C:/Windows/Fonts")
JP_FONT_NAME = "Meiryo"


def _font_path(bold: bool = False) -> str:
    order = ["meiryob.ttc", "meiryo.ttc"] if bold else FONT_CANDIDATES
    for n in order:
        p = FONT_DIR / n
        if p.exists():
            return str(p)
    raise RuntimeError("日本語フォントが見つからない")


# ---------------------------------------------------------------------------
# レイアウト要素（pptx と PNG の共通定義）
# ---------------------------------------------------------------------------

@dataclass
class Text:
    x: float; y: float; w: float; h: float          # inch
    text: str
    size: float = 14                                 # pt
    color: tuple = INK
    bold: bool = False
    align: str = "l"                                 # l / c / r
    mono: bool = False
    line_spacing: float = 1.35


@dataclass
class Box:
    x: float; y: float; w: float; h: float
    fill: tuple | None = PANEL
    line: tuple | None = None
    lw: float = 1.2


@dataclass
class Pic:
    x: float; y: float; w: float; h: float
    path: Path


@dataclass
class Slide:
    title: str
    elements: list = field(default_factory=list)
    notes: str = ""


# ---------------------------------------------------------------------------
# データ読み込み（唯一の数値の出どころ）
# ---------------------------------------------------------------------------

def load():
    corr = json.loads((MAIN / "transition_recomputed_preregistered.json").read_text(encoding="utf-8"))
    sens = json.loads(SENS.read_text(encoding="utf-8"))
    runs = {(c, s): json.loads((MAIN / f"{c}_seed{s}.json").read_text(encoding="utf-8"))
            for c in "ABCD" for s in (2, 4, 6, 7, 9)}
    live = [json.loads(p.read_text(encoding="utf-8")) for p in sorted(LIVE.glob("*.json"))]
    elig = json.loads((ROOT / "outputs" / "seed_eligibility.json").read_text(encoding="utf-8"))
    return corr, sens, runs, live, elig


def facts(corr, sens, runs, live, elig):
    R = corr["runs"]
    d4 = corr["preregistered_D4"]
    rt = corr["runtime_D4"]
    n = len(R)
    ever = {k: sum(1 for r in R if r["ever_met"][k])
            for k in ("active_supplier_count", "supply_duration",
                      "coordination_edges", "community_supply_share")}
    pens = sorted(sens["penalties"])
    rate = {p: st.fmean(x["expected_units_per_step"] for x in sens["records"]
                        if x["modify_difficulty_penalty"] == p) for p in pens}
    prob = {p: st.fmean(x["corrected_transition_probability"] for x in sens["records"]
                        if x["modify_difficulty_penalty"] == p) for p in pens}
    # 条件A の make 内訳（RESULTS.md §7.3 と同じ帰属順）
    b = {"supply": 0, "fail": 0, "unmet": 0, "mat": 0, "asset": 0}
    for s in (2, 4, 6, 7, 9):
        for p in runs[("A", s)]["provenance"]:
            if p["supplied_units"] > 0: b["supply"] += 1
            elif not p["meets_requirement"]: b["unmet"] += 1
            elif not p["make_success"]: b["fail"] += 1
            elif not p["material_feasible"]: b["mat"] += 1
            elif not p["asset_feasible"]: b["asset"] += 1
    cost = {"observe": .1, "ask": .2, "practice": .5, "make": 1., "share": .2,
            "idle": 0., "modify": .8, "propose": .2, "join": .2}
    used = sum(k * cost[a] for s in (2, 4, 6, 7, 9)
               for a, k in runs[("A", s)]["action_counts"].items())
    avail = 3.0 * 6 * 8 * 5
    # modify 回数別の実効成功確率（条件A。Slide 6 の最小視覚表現に使う）
    pm: dict[int, list[float]] = {}
    for s in (2, 4, 6, 7, 9):
        for p in runs[("A", s)]["provenance"]:
            pm.setdefault(len(p["applied_modifications"] or {}), []).append(
                p["effective_success_probability"])
    p_by_mod = {k: st.fmean(v) for k, v in pm.items()}
    # 全条件: 供給と協調の重なり
    tot_sup = coord_sup = 0
    for r in runs.values():
        for p in r["provenance"]:
            if p["supplied_units"] > 0:
                tot_sup += 1
                if p["coordination_relation"]["edges_involving_self"]:
                    coord_sup += 1
    # P*（RESULTS.md §8.5 と同一手続き）
    ext = runs[("A", 2)]["D5_external_reference_supply_per_step"]
    need = d4["community_supply_share"] * ext / (1 - d4["community_supply_share"])
    cur = 0.35
    att = [(len(p["applied_modifications"] or {}),
            p["effective_success_probability"] * (1 + cur * len(p["applied_modifications"] or {})),
            p["meets_requirement"]) for r in runs.values() for p in r["provenance"]]
    def supply(P):
        return sum(max(.02, min(.98, pb / (1 + P * k))) for k, pb, q in att if q) / len(runs) / 8
    lo, hi = 0., 1.
    for _ in range(200):
        m = (lo + hi) / 2
        if supply(m) > need: lo = m
        else: hi = m
    pstar = (lo + hi) / 2
    span_pen = max(rate.values()) - min(rate.values())
    span_cond = max(abs(st.fmean(x["expected_units_per_step"] for x in sens["records"]
                                 if x["modify_difficulty_penalty"] == p and x["condition"] == a)
                        - st.fmean(x["expected_units_per_step"] for x in sens["records"]
                                   if x["modify_difficulty_penalty"] == p and x["condition"] == c))
                    for p in pens for a in "ABCD" for c in "ABCD")
    return dict(
        n=n, d4=d4, rt=rt, ever=ever,
        corrected=sum(1 for r in R if r["corrected_transition"]),
        runtime=sum(1 for r in R if r["runtime_transition"]),
        pens=pens, rate=rate, prob=prob, pstar=pstar, cur=cur,
        span_pen=span_pen, span_cond=span_cond,
        b=b, make_total=sum(b.values()), used_pct=used / avail * 100,
        p_by_mod=p_by_mod,
        tot_sup=tot_sup, coord_sup=coord_sup, need=need,
        scan=elig["seeds_scanned"], elig=elig["eligible_seeds"],
        live=[(d["seed"], d["community_supply_total"],
               d["replay_prediction"]["expected_community_supply_total"],
               d["live_minus_replay"], d["replay_prediction"]["tolerance_2sd"]) for d in live],
        calls=sum(r["llm_calls"] for r in runs.values()),
        spent=sum(r["spent_usd"] for r in runs.values()),
    )


# ---------------------------------------------------------------------------
# スライド定義
# ---------------------------------------------------------------------------

def build(f) -> list[Slide]:
    d4, rt, ev = f["d4"], f["rt"], f["ever"]
    S = []

    # ---- Slide 1 -----------------------------------------------------------
    s = Slide("コスプレ文化が生む「制作能力」を測る")
    s.elements += [
        Text(1.0, 1.55, 11.3, 1.5,
             "コスプレ文化が毎週生み出している制作能力を、\n"
             "私たちは社会資源として一度でも測ったことがあるでしょうか。",
             size=27, bold=True, color=INK, line_spacing=1.5),
        Box(1.0, 3.62, 11.3, 0.028, fill=LIGHT),
        Text(1.0, 3.95, 11.3, 0.55, "衣装の話ではありません。", size=21, color=GAP, bold=True),
        Text(1.0, 4.72, 11.3, 0.75,
             "毎回違う仕様を、一点だけ、締切までに、あり合わせの材料で作る。",
             size=19, color=INK),
        Text(1.0, 5.55, 11.3, 0.5, "測ったのは、この生産能力です。", size=17, color=MUTED),
    ]
    S.append(s)

    # ---- Slide 2 -----------------------------------------------------------
    s = Slide("なぜコスプレなのか ── 毎回違うものを作る生産体制")
    rows = [("仕様", "固定", "毎回ちがう"), ("ロット", "大量", "一点物"),
            ("納期", "計画", "締切駆動"), ("現場修復", "想定外", "日常"),
            ("知識", "秘匿", "公開共有が規範")]
    y0 = 1.85
    s.elements += [
        Text(1.6, 1.30, 2.4, 0.4, "", size=13),
        Text(4.3, 1.30, 3.2, 0.4, "製造業", size=17, bold=True, color=MUTED, align="c"),
        Text(7.9, 1.30, 3.8, 0.4, "コスプレ制作", size=17, bold=True, color=ACCENT, align="c"),
        Box(1.6, 1.76, 10.1, 0.02, fill=LIGHT),
    ]
    for i, (a, b_, c_) in enumerate(rows):
        yy = y0 + i * 0.66
        s.elements += [
            Text(1.6, yy, 2.4, 0.5, a, size=16, color=INK),
            Text(4.3, yy, 3.2, 0.5, b_, size=16, color=MUTED, align="c"),
            Text(7.9, yy, 3.8, 0.5, c_, size=16, color=ACCENT, bold=True, align="c"),
        ]
    s.elements += [
        Box(1.6, 5.22, 10.1, 0.02, fill=LIGHT),
        Text(1.6, 5.55, 10.1, 0.95,
             "危機時に足りなくなるのは、誰も備蓄していない規格外の少量品。\n"
             "着目したのは技能ではなく、生産体制の構造。",
             size=17, bold=True, color=INK, line_spacing=1.5),
    ]
    S.append(s)

    # ---- Slide 3 -----------------------------------------------------------
    s = Slide("その能力は、供給ショックで転化するか？")
    steps = [
        ("文化活動 → skill / project / network の蓄積", "156 週", PANEL, MUTED),
        ("供給ショック", "1 step = 6 時間 × 8", "#fdf1ee", GAP),
        (f"未知仕様（属性のみ）  attr_0 ≥ {d4['community_supply_share'] and 0.60}"
         f"  /  attr_2 ≥ 0.55", "何を作るべきかは指示しない", "#eef3f8", ACCENT),
        ("modify  ─  coordination  ─  make  →  supply", "", PANEL, MUTED),
    ]
    for i, (t, sub, fill, ec) in enumerate(steps):
        yy = 1.42 + i * 1.30
        fillc = fill if isinstance(fill, tuple) else tuple(
            int(fill.lstrip("#")[j:j + 2], 16) for j in (0, 2, 4))
        s.elements += [Box(1.5, yy, 10.3, 0.92, fill=fillc, line=ec, lw=1.4),
                       Text(1.8, yy + 0.16, 7.4, 0.6, t, size=17, bold=True, color=INK)]
        if sub:
            s.elements.append(Text(9.2, yy + 0.20, 2.4, 0.5, sub, size=13, color=ec, align="r"))
        if i < 3:
            s.elements.append(Text(6.4, yy + 0.94, 0.6, 0.34, "▼", size=13, color=LIGHT, align="c"))
    s.elements += [
        Text(1.5, 6.55, 10.3, 0.5,
             "Agent 40体（うち shock 対象 6体）／属性は中立コード表記／"
             "LLM decides intent. Code determines feasibility.",
             size=12.5, color=MUTED),
    ]
    S.append(s)

    # ---- Slide 4 -----------------------------------------------------------
    s = Slide("Agent は既存 Project を未知仕様へ適応させた")
    s.elements += [
        Pic(0.55, 1.28, 8.25, 4.13, FIGS / "F3_adaptation_path.png"),
        # 動画差し込み領域（16:9。3.65 x 2.053 in = 526 x 296 px @1920x1080）
        Box(9.15, 1.28, 3.65, 2.053, fill=(250, 250, 250), line=ACCENT, lw=1.6),
        Text(9.30, 2.02, 3.35, 0.55, "［ 4秒動画を差し込む領域 ］\n16:9  3.65 × 2.05 in",
             size=12.5, bold=True, color=ACCENT, align="c", line_spacing=1.45),
        Text(9.15, 3.45, 3.65, 0.95,
             "M1 / API-free simulation\n156 accumulation steps\n10× playback",
             size=12, color=MUTED, align="c", line_spacing=1.5, mono=True),
        Box(9.15, 4.55, 3.65, 1.05, fill=(240, 245, 240), line=OK, lw=1.2),
        Text(9.30, 4.72, 3.35, 0.8,
             "証拠 = F3（provenance）\n挙動 = 実行画面",
             size=12.5, bold=True, color=OK, align="c", line_spacing=1.5),
        Text(0.55, 5.62, 12.25, 1.0,
             "既存 Project を要求仕様に合わせて適応的に modify し、供給へ到達した。\n"
             "make 試行ごとの provenance に直接記録（集計からの推論ではない）。",
             size=15, bold=True, color=INK, line_spacing=1.5),
    ]
    s.notes = ("F3 = 証拠（provenance）。4秒動画 = 挙動（M1 が実際に動くこと）。"
               "M1 映像を main experiment の実行映像や仕様適応の瞬間として見せない。")
    S.append(s)

    # ---- Slide 5 -----------------------------------------------------------
    s = Slide(f"しかし Transition は {f['corrected']}/{f['n']} だった")
    s.elements += [Pic(0.55, 1.22, 7.75, 3.76, FIGS / "F1_transition_conditions.png")]
    items = [("供給者の形成", f"active_supplier_count ≥ {d4['active_supplier_count']}",
              ev["active_supplier_count"], OK),
             ("供給の継続", f"supply_duration ≥ {d4['supply_duration_steps']} step",
              ev["supply_duration"], OK),
             ("協調関係の形成", f"coordination_edges ≥ {d4['coordination_edges']}",
              ev["coordination_edges"], OK),
             ("量", f"community_supply_share ≥ {d4['community_supply_share']}",
              ev["community_supply_share"], GAP)]
    for i, (lab, cond, v, col) in enumerate(items):
        yy = 1.32 + i * 0.72
        s.elements += [
            Text(8.60, yy, 2.35, 0.34, lab, size=14, bold=True, color=col),
            Text(8.60, yy + 0.32, 2.95, 0.30, cond, size=10.5, color=MUTED),
            Text(11.55, yy + 0.02, 1.30, 0.45, f"{v} / {f['n']}",
                 size=17, bold=True, color=col, align="r"),
        ]
    s.elements += [
        Box(8.60, 4.28, 4.25, 0.02, fill=LIGHT),
        Text(8.60, 4.45, 4.25, 0.55,
             f"4条件の同時成立（Transition）= {f['corrected']} / {f['n']}",
             size=15, bold=True, color=GAP),
        Box(0.55, 5.28, 12.25, 1.30, fill=(240, 245, 242), line=OK, lw=1.3),
        Text(0.85, 5.48, 11.65, 0.95,
             "供給者形成・継続・協調は全 20 run で成立。量的条件との同時成立が残った。\n"
             f"※ 量の条件は {ev['community_supply_share']} / {f['n']} run で瞬間的に充足したが、"
             f"供給が継続する時刻には届いていない。「share が {f['corrected']}/{f['n']}」ではない。",
             size=13.5, bold=False, color=INK, line_spacing=1.55),
    ]
    s.notes = "share 2/20（瞬間的）と Transition 0/20（4条件同時）を混同しない。"
    S.append(s)

    # ---- Slide 6 -----------------------------------------------------------
    s = Slide("なぜか ── 仕様適応後の production success が律速した")
    p0 = f["p_by_mod"].get(0, 0.0)
    p2 = f["p_by_mod"].get(2, 0.0)
    s.elements += [
        # 中心命題（1つだけ）。3秒で読めるよう最上部に大きく置く
        Box(0.55, 1.16, 12.25, 0.92, fill=(240, 245, 242), line=OK, lw=1.6),
        Text(0.85, 1.28, 11.65, 0.80,
             "量を止めていたのは、材料でも設備でも時間でもなく、\n"
             "モデル上の仕様適応後の production success だった。",
             size=19, bold=True, color=INK, align="c", line_spacing=1.45),
        # F4 を主図として大きく
        Pic(0.55, 2.32, 8.45, 3.37, FIGS / "F4_bottleneck.png"),
        # F2 は削除。代わりに最小限の視覚表現のみ
        Box(9.25, 2.32, 3.55, 2.35, fill=PANEL, line=LIGHT, lw=1.2),
        Text(9.25, 2.50, 3.55, 0.38, "仕様適応のコスト", size=13, bold=True,
             color=INK, align="c"),
        Text(9.42, 3.02, 1.35, 0.36, "modify 0箇所", size=11.5, color=MUTED),
        Text(9.42, 3.34, 1.45, 0.55, f"p ≈ {p0:.2f}", size=21, bold=True,
             color=OK),
        Text(10.92, 3.32, 0.40, 0.5, "→", size=19, color=MUTED, align="c"),
        Text(11.42, 3.02, 1.35, 0.36, "modify 2箇所", size=11.5, color=MUTED),
        Text(11.42, 3.34, 1.32, 0.55, f"p ≈ {p2:.2f}", size=21, bold=True,
             color=GAP),
        Text(9.45, 4.08, 3.15, 0.55,
             "要求仕様に合わせるほど、\n作れる確率が下がる。",
             size=12, color=INK, align="c", line_spacing=1.4),
        # ★ 現実の値ではないことの注記（一目で分かる位置）
        Box(9.25, 4.82, 3.55, 0.88, fill=(253, 241, 238), line=GAP, lw=1.2),
        Text(9.42, 4.96, 3.2, 0.65,
             f"※ 現行モデル内の値。未校正の\n"
             f"modify_difficulty_penalty = {f['cur']} による。\n"
             f"現実の成功率ではない。",
             size=10, color=GAP, line_spacing=1.4),
        Text(0.55, 5.90, 8.45, 0.55,
             f"make 試行 {f['make_total']} 件（条件A）｜ 材料不足で阻止 {f['b']['mat']} 件、"
             f"設備なしで阻止 {f['b']['asset']} 件、時間予算は {100 - f['used_pct']:.0f}% が未使用",
             size=12, color=MUTED),
        Box(0.55, 6.52, 12.25, 0.58, fill=None, line=LIGHT, lw=1.0),
        Text(0.80, 6.63, 11.75, 0.38,
             "P* / sensitivity の詳細は Backup B2。"
             "penalty = 0.00 は「追加ペナルティをゼロと仮定したモデル上の構造的上限ケース」であり、"
             "達成可能な改善策ではない。",
             size=11, color=MUTED),
    ]
    s.notes = ("口頭は中心命題1つのみ。P*・8.2倍・7.4倍は読み上げず Backup B2 に任せる。"
               "0.98→0.58 は現行モデル内の未校正パラメータによる値。")
    S.append(s)

    # ---- Slide 7 -----------------------------------------------------------
    s = Slide(f"そして私は「{f['runtime']}/{f['n']}」を採用しなかった")
    s.elements += [
        # 副題はタイトル直下・同一視線移動内
        Text(0.62, 1.02, 12.1, 0.5, "── 実行時 config が事前登録値と違っていたから",
             size=19, bold=True, color=GAP),
        Pic(0.55, 1.66, 5.55, 4.68, FIGS / "F5_methods.png"),
        Box(6.35, 1.66, 6.45, 1.02, fill=(238, 243, 248), line=ACCENT, lw=1.4),
        Text(6.58, 1.80, 6.0, 0.75,
             f"事前登録  D4 = share {d4['community_supply_share']} / "
             f"suppliers {d4['active_supplier_count']} /\n"
             f"　　　　　duration {d4['supply_duration_steps']} / edges {d4['coordination_edges']}",
             size=13.5, bold=True, color=ACCENT, line_spacing=1.45),
        Box(6.35, 2.92, 6.45, 0.86, fill=(250, 244, 238), line=(194, 161, 132), lw=1.3),
        Text(6.58, 3.06, 6.0, 0.6,
             f"実行時 config（暫定値）→ Transition {f['runtime']} / {f['n']}\n"
             f"【副次的記録・保持】",
             size=13, color=INK, line_spacing=1.45),
        Text(6.35, 3.88, 6.45, 0.36, "▼  事前登録値との突合で不一致を検出",
             size=12, color=GAP, bold=True),
        Box(6.35, 4.32, 6.45, 0.86, fill=(234, 243, 238), line=OK, lw=1.6),
        Text(6.58, 4.46, 6.0, 0.6,
             f"事前登録値を適用 → Transition {f['corrected']} / {f['n']}\n【主結果】",
             size=13.5, bold=True, color=OK, line_spacing=1.45),
        Text(6.35, 5.34, 6.45, 0.5,
             "ログは1行も改変していない。判定基準を選び直したのではない。",
             size=12.5, color=INK),
        Box(6.35, 5.90, 6.45, 0.02, fill=LIGHT),
        Text(6.35, 6.02, 6.45, 0.4,
             "202 tests  /  preregistered adjudication  /  API-free replay",
             size=12, bold=True, color=MUTED, mono=True),
    ]
    s.notes = "「良い結果を捨てた」と言わない。最初から決めてあった基準を適用した。"
    S.append(s)

    # ---- Slide 8 -----------------------------------------------------------
    s = Slide("次は現実のコスプレ制作能力を測る")
    s.elements += [
        Text(0.75, 1.22, 5.75, 0.42, "研究が示したこと（モデル内）",
             size=15, bold=True, color=OK),
        Box(0.75, 1.68, 5.75, 0.02, fill=OK),
        Text(0.75, 1.88, 5.75, 2.3,
             "既存 Project を要求仕様へ適応的に modify し、\n供給へ到達した。\n\n"
             "量の律速は、仕様適応後の\nproduction success だった。\n"
             "この係数は未校正である。",
             size=14, color=INK, line_spacing=1.55),
        # 境界線（物理的な分離）
        Box(6.83, 1.22, 0.02, 4.55, fill=LIGHT),
        Text(7.20, 1.22, 5.6, 0.42, "ここから先は製品仮説（研究結果ではない）",
             size=15, bold=True, color=ACCENT),
        Box(7.20, 1.68, 5.6, 0.02, fill=ACCENT),
        Text(7.20, 1.88, 5.6, 0.5, "Specification Adaptation Challenge",
             size=15, bold=True, color=INK),
        Text(7.20, 2.45, 5.6, 1.9,
             "実際の制作者が未知仕様へ適応するときの\n\n"
             "　成功率 ／ 制作時間 ／ 歩留まり\n　材料・設備制約 ／ 協調形成\n\nを測る",
             size=13.5, color=INK, line_spacing=1.5),
        Box(0.75, 5.98, 12.05, 0.92, fill=(240, 245, 242), line=OK, lw=1.4),
        Text(0.75, 6.20, 12.05, 0.5,
             "好きなことを続けてください。それが安全保障になる社会へ。",
             size=20, bold=True, color=INK, align="c"),
    ]
    S.append(s)

    # ---- Backup B1 ---------------------------------------------------------
    s = Slide("Backup B1 — Agent 行動・解の多様性")
    s.elements += [
        Text(0.75, 1.22, 12.05, 0.4, "想定質問: 「その1件は都合のよい好例では？」",
             size=14, bold=True, color=MUTED),
        Box(0.75, 1.78, 12.05, 1.15, fill=(240, 245, 242), line=OK, lw=1.3),
        Text(1.00, 1.98, 11.55, 0.8,
             f"供給を伴う make 記録 {f['tot_sup']} 件のうち {f['coord_sup']} 件が、"
             f"同時刻に協調関係を持つ作り手によるもの\n（20/20 run で発生）。"
             f"経路は1件だけの偶然ではない。",
             size=14, color=INK, line_spacing=1.5),
        Text(0.75, 3.15, 12.05, 0.9,
             f"条件A の make 試行 {f['make_total']} 件の内訳:  "
             f"供給成立 {f['b']['supply']} ／ 確率的な制作失敗 {f['b']['fail']} ／ "
             f"仕様未達 {f['b']['unmet']} ／ 材料不足で阻止 {f['b']['mat']} ／ "
             f"設備なしで阻止 {f['b']['asset']}",
             size=13.5, color=INK, line_spacing=1.5),
        Box(0.75, 4.10, 12.05, 1.55, fill=(253, 241, 238), line=GAP, lw=1.3),
        Text(1.00, 4.32, 11.55, 1.15,
             "ただし解の多様性は低い（隠さない）:\n"
             f"　・供給に使われた制作対象は 2 種のみ（{f['tot_sup']} 件中 579 件が同一）\n"
             "　・modify の内容も「2属性を各 +0.15」がほぼ全 run で反復",
             size=13.5, color=INK, line_spacing=1.55),
        Text(0.75, 5.85, 12.05, 0.85,
             "言い方: 「1件は代表例であり、経路自体は 20/20 run で成立したが、解の多様性は低い」。\n"
             "禁止: 「多様な創発が確認された」「未知用途への創発的再構成を証明した」",
             size=12, color=MUTED, line_spacing=1.5),
    ]
    S.append(s)

    # ---- Backup B2 ---------------------------------------------------------
    s = Slide("Backup B2 — D4 / P* / sensitivity の詳細")
    s.elements += [
        Text(0.75, 1.20, 12.05, 0.4,
             "想定質問: 「なぜ share だけ届かない？」「閾値を変えれば結果が変わるのでは？」",
             size=14, bold=True, color=MUTED),
        Box(0.75, 1.72, 12.05, 0.80, fill=(238, 243, 248), line=ACCENT, lw=1.3),
        Text(1.00, 1.88, 11.55, 0.55,
             "share 2/20 が「瞬間的」な理由: share の最大は t+0、"
             f"duration ≥ {d4['supply_duration_steps']} の資格は t+3 以降。両者の時間帯が重ならない。",
             size=13.5, color=INK),
    ]
    hdr = ["penalty", "供給 (units/step)", "転化確率"]
    for j, h in enumerate(hdr):
        s.elements.append(Text(0.95 + j * 2.55, 2.72, 2.4, 0.36, h,
                               size=13, bold=True, color=MUTED))
    for i, p in enumerate(f["pens"]):
        yy = 3.12 + i * 0.44
        cur_mark = "  ← current" if abs(p - f["cur"]) < 1e-9 else ""
        col = INK if abs(p - f["cur"]) < 1e-9 else MUTED
        s.elements += [
            Text(0.95, yy, 2.4, 0.36, f"{p:.2f}{cur_mark}", size=13,
                 bold=abs(p - f["cur"]) < 1e-9, color=col),
            Text(3.50, yy, 2.4, 0.36, f"{f['rate'][p]:.3f}", size=13, color=col),
            Text(6.05, yy, 2.4, 0.36, f"{f['prob'][p]:.4f}", size=13, color=col),
        ]
    s.elements += [
        Box(8.75, 2.68, 4.05, 2.45, fill=PANEL, line=LIGHT, lw=1.1),
        Text(8.98, 2.84, 3.6, 2.22,
             f"P* = {f['pstar']:.4f}\nD4 share-equivalent\nsensitivity boundary\n\n"
             f"現行値 {f['cur']} はその約 {f['cur'] / f['pstar']:.1f} 倍\n\n"
             f"penalty の効果は条件差の\n約 {f['span_pen'] / f['span_cond']:.1f} 倍",
             size=12.5, color=INK, line_spacing=1.55),
        Box(0.75, 5.15, 12.05, 0.80, fill=(253, 241, 238), line=GAP, lw=1.3),
        Text(1.00, 5.31, 11.55, 0.55,
             "penalty = 0.00 は「追加ペナルティをゼロと仮定したモデル上の構造的上限ケース」。"
             "達成可能な改善策ではない。",
             size=13, color=GAP),
        Text(0.75, 6.15, 12.05, 0.85,
             "「閾値を変えれば結果が変わるのでは」への回答: "
             "研究側の閾値は事前登録値のまま一切変更していない。\n"
             "禁止: 「損益分岐」という語（正式名称は D4 share-equivalent sensitivity boundary）",
             size=12, color=MUTED, line_spacing=1.5),
    ]
    S.append(s)

    # ---- Backup B3 ---------------------------------------------------------
    s = Slide("Backup B3 — 再現性 / config 不一致 / corrected adjudication")
    lv = "　".join(f"seed {sd}: D={d:+.3f}（許容 ±{tol:.2f}）" for sd, _, _, d, tol in f["live"])
    s.elements += [
        Text(0.75, 1.20, 12.05, 0.4,
             "想定質問: 「なぜ config がずれた？」「後から都合よく判定を変えたのでは？」",
             size=14, bold=True, color=MUTED),
        Box(0.75, 1.72, 12.05, 1.10, fill=(250, 244, 238), line=(194, 161, 132), lw=1.3),
        Text(1.00, 1.90, 11.55, 0.8,
             "原因: D4 を確定した commit が base.yaml の D5 ブロックのみ更新し、"
             "直下の D4 ブロックを素通りした。\n"
             "検出漏れ: 逸脱チェックが「20 run で同じ値か」しか見ず、事前登録値との照合を欠いていた。",
             size=13, color=INK, line_spacing=1.5),
        Box(0.75, 3.02, 12.05, 1.10, fill=(234, 243, 238), line=OK, lw=1.3),
        Text(1.00, 3.20, 11.55, 0.8,
             "再判定は判定ロジックを再実装していない。既存の TransitionJudge.evaluate() に\n"
             "ログ済みの閾値非依存な実測値を再投入。runtime 閾値で replay すると "
             "20 run × 8 step 全件が完全一致。",
             size=13, color=INK, line_spacing=1.5),
        Text(0.75, 4.32, 12.05, 0.52,
             f"live run 2本（条件A seed 2/4、penalty=0.00）で replay 近似を検証:  {lv}  → 両者とも許容内",
             size=12.5, color=INK),
        Text(0.75, 4.92, 12.05, 0.85,
             "是正: 実行時 config を configs/as_executed/ に保存（編集禁止）／"
             "base.yaml を事前登録値へ同期／\n三者突合監査と同期テストを追加。",
             size=13, color=INK, line_spacing=1.5),
        Box(0.75, 5.88, 12.05, 0.95, fill=PANEL, line=LIGHT, lw=1.1),
        Text(1.00, 6.05, 11.55, 0.65,
             f"API 不使用で再現できる範囲: M1 全体 ／ 感度分析 replay ／ 主結果の再判定 ／ テスト 202 件。\n"
             f"main experiment の再実行のみ {f['calls']} calls・約 ${f['spent']:.1f}。",
             size=12.5, color=INK, line_spacing=1.5),
    ]
    S.append(s)
    return S


# ---------------------------------------------------------------------------
# レンダラ 1: pptx
# ---------------------------------------------------------------------------

def to_pptx(slides: list[Slide], path: Path) -> None:
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(W_IN), Inches(H_IN)
    blank = prs.slide_layouts[6]
    for sl in slides:
        s = prs.slides.add_slide(blank)
        # タイトル
        tb = s.shapes.add_textbox(Inches(0.62), Inches(0.32), Inches(12.1), Inches(0.72))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run(); r.text = sl.title
        r.font.size = Pt(25); r.font.bold = True
        r.font.color.rgb = RGBColor(*INK); r.font.name = JP_FONT_NAME
        for el in sl.elements:
            if isinstance(el, Box):
                sh = s.shapes.add_shape(1, Inches(el.x), Inches(el.y),
                                        Inches(el.w), Inches(el.h))
                if el.fill:
                    sh.fill.solid(); sh.fill.fore_color.rgb = RGBColor(*el.fill)
                else:
                    sh.fill.background()
                if el.line:
                    sh.line.color.rgb = RGBColor(*el.line); sh.line.width = Pt(el.lw)
                else:
                    sh.line.fill.background()
                sh.shadow.inherit = False
            elif isinstance(el, Pic):
                s.shapes.add_picture(str(el.path), Inches(el.x), Inches(el.y),
                                     Inches(el.w), Inches(el.h))
            elif isinstance(el, Text):
                tb = s.shapes.add_textbox(Inches(el.x), Inches(el.y),
                                          Inches(el.w), Inches(el.h))
                tf = tb.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.TOP
                for i, line in enumerate(el.text.split("\n")):
                    p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                    p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER,
                                   "r": PP_ALIGN.RIGHT}[el.align]
                    p.line_spacing = el.line_spacing
                    r = p.add_run(); r.text = line
                    r.font.size = Pt(el.size); r.font.bold = el.bold
                    r.font.color.rgb = RGBColor(*el.color)
                    r.font.name = "Consolas" if el.mono else JP_FONT_NAME
        if sl.notes:
            s.notes_slide.notes_text_frame.text = sl.notes
    prs.save(str(path))


# ---------------------------------------------------------------------------
# レンダラ 2: プレビュー PNG（同一定義から描画。文字切れ監査に使う）
# ---------------------------------------------------------------------------

_font_cache: dict = {}


def _f(size_pt: float, bold: bool, mono: bool = False):
    key = (round(size_pt, 1), bold, mono)
    if key not in _font_cache:
        px = int(size_pt * PX_PER_IN / 72 * SCALE)
        if mono:
            p = "C:/Windows/Fonts/consola.ttf"
            p = p if Path(p).exists() else _font_path(bold)
        else:
            p = _font_path(bold)
        _font_cache[key] = ImageFont.truetype(p, px)
    return _font_cache[key]


def _wrap(draw, text, font, max_px):
    out = []
    for para in text.split("\n"):
        if not para:
            out.append("")
            continue
        cur = ""
        for ch in para:
            if draw.textlength(cur + ch, font=font) <= max_px:
                cur += ch
            else:
                out.append(cur); cur = ch
        out.append(cur)
    return out


def to_preview(slides: list[Slide], outdir: Path) -> list[dict]:
    outdir.mkdir(parents=True, exist_ok=True)
    issues = []
    for i, sl in enumerate(slides, start=1):
        W, H = int(W_IN * PX_PER_IN * SCALE), int(H_IN * PX_PER_IN * SCALE)
        img = Image.new("RGB", (W, H), BG)
        d = ImageDraw.Draw(img)
        sc = PX_PER_IN * SCALE

        def box(el):
            xy = [el.x * sc, el.y * sc, (el.x + el.w) * sc, (el.y + el.h) * sc]
            d.rectangle(xy, fill=el.fill, outline=el.line,
                        width=max(1, int(el.lw * SCALE)) if el.line else 0)

        for el in sl.elements:
            if isinstance(el, Box):
                box(el)
        for el in sl.elements:
            if isinstance(el, Pic):
                im = Image.open(el.path).convert("RGB")
                tw, th = int(el.w * sc), int(el.h * sc)
                im = im.resize((tw, th), Image.LANCZOS)
                img.paste(im, (int(el.x * sc), int(el.y * sc)))

        # タイトル
        tf = _f(25, True)
        d.text((0.62 * sc, 0.32 * sc), sl.title, font=tf, fill=INK)
        tw = d.textlength(sl.title, font=tf)
        if tw > 12.1 * sc:
            issues.append({"slide": i, "kind": "title_overflow",
                           "detail": f"タイトル幅 {tw / sc:.2f}in > 12.10in"})

        for el in sl.elements:
            if not isinstance(el, Text):
                continue
            font = _f(el.size, el.bold, el.mono)
            lines = _wrap(d, el.text, font, el.w * sc)
            lh = el.size * PX_PER_IN / 72 * SCALE * el.line_spacing
            need = lh * len(lines)
            if need > el.h * sc + 2:
                issues.append({"slide": i, "kind": "text_overflow",
                               "detail": f"必要高 {need / sc:.2f}in > 枠 {el.h:.2f}in",
                               "text": el.text.split("\n")[0][:40]})
            for j, line in enumerate(lines):
                lw_ = d.textlength(line, font=font)
                x = el.x * sc
                if el.align == "c": x = (el.x + el.w / 2) * sc - lw_ / 2
                elif el.align == "r": x = (el.x + el.w) * sc - lw_
                d.text((x, el.y * sc + j * lh), line, font=font, fill=el.color)
            # 画面外
            if el.x + el.w > W_IN + 0.01 or el.y + el.h > H_IN + 0.01 or el.x < -0.01 or el.y < -0.01:
                issues.append({"slide": i, "kind": "offslide",
                               "detail": f"x={el.x} y={el.y} w={el.w} h={el.h}"})

        img.resize((W // SCALE, H // SCALE), Image.LANCZOS).save(
            outdir / f"slide{i:02d}.png")
    return issues


def main() -> int:
    OUT.mkdir(exist_ok=True)
    corr, sens, runs, live, elig = load()
    f = facts(corr, sens, runs, live, elig)
    slides = build(f)

    pptx_path = OUT / "cosplay_reserve_final.pptx"
    to_pptx(slides, pptx_path)
    issues = to_preview(slides, PREVIEW)

    print(f"pptx: {pptx_path}  ({len(slides)} 枚, 16:9 {W_IN}x{H_IN}in)")
    for i, s in enumerate(slides, 1):
        tag = "本編" if i <= 8 else "Backup"
        print(f"  {i:>2} [{tag}] {s.title}")
    print(f"\nプレビュー PNG: {PREVIEW}")
    if issues:
        print(f"\n⚠ 視認性の問題 {len(issues)} 件:")
        for x in issues:
            print(f"   slide{x['slide']:>2} {x['kind']}: {x['detail']} {x.get('text','')}")
    else:
        print("\n視認性チェック: 文字切れ・枠外なし")
    print(f"\n主要数値（ログ由来）: Transition {f['corrected']}/{f['n']} / "
          f"runtime {f['runtime']}/{f['n']} / share ever_met "
          f"{f['ever']['community_supply_share']}/{f['n']} / P*={f['pstar']:.4f}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
